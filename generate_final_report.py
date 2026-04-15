"""Generate final comprehensive WBR report with WoW analysis + charts."""
import json
import os
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches

OUTPUT = "/home/mousonko/.workspace/wbr-bridge-agent/output"

with open(f"{OUTPUT}/wbr_analysis.json") as f:
    analysis = json.load(f)
with open(f"{OUTPUT}/wow_analysis.json") as f:
    wow = json.load(f)


def img(name):
    return os.path.join(OUTPUT, name)


def add_img(doc, name, width=6.2):
    p = img(name)
    if os.path.exists(p):
        doc.add_picture(p, width=Inches(width))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER


# ============ WORD REPORT ============
doc = Document()
style = doc.styles["Normal"]
style.font.name = "Calibri"
style.font.size = Pt(10)

doc.add_heading("EU LM WBR — CRET Scan Compliance Bridge", level=0).alignment = WD_ALIGN_PARAGRAPH.CENTER
doc.add_paragraph("Week 14 | Comprehensive Analysis with Week-over-Week Trends",
                  style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER

# Key Metrics
doc.add_heading("Key Metrics", level=1)
t = doc.add_table(rows=2, cols=6)
t.style = "Light Grid Accent 1"
t.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h in enumerate(["Flagged Sites", "With Bridge", "No Bridge", "Deteriorating", "New Flags", "Resolved"]):
    t.rows[0].cells[i].text = h
t.rows[1].cells[0].text = str(analysis.get("total_flagged_sites", ""))
t.rows[1].cells[1].text = str(analysis.get("sites_with_bridges", ""))
t.rows[1].cells[2].text = str(analysis.get("sites_without_bridges", ""))
t.rows[1].cells[3].text = str(len(wow.get("deteriorating", [])))
t.rows[1].cells[4].text = str(len(wow.get("new_flags", [])))
t.rows[1].cells[5].text = str(len(wow.get("resolved", [])))

# Executive Summary
doc.add_heading("Executive Summary", level=1)
doc.add_paragraph(analysis.get("executive_summary", ""))

# WoW Summary
doc.add_heading("Week-over-Week Summary", level=1)
doc.add_paragraph(wow.get("summary", ""))

# WoW Trending Chart
doc.add_heading("Compliance Trend — WK-12 → WK-14", level=1)
add_img(doc, "chart_wow_trending.png")

# Flagged Sites per Week
doc.add_heading("Flagged Sites per Week by Country", level=1)
add_img(doc, "chart_wow_site_count.png")

# Deteriorating Sites
doc.add_heading("Deteriorating Sites: WK-13 → WK-14", level=1)
add_img(doc, "chart_deteriorating.png", 6.5)

det = wow.get("deteriorating", [])
if det:
    t = doc.add_table(rows=1, cols=5)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Site", "MP", "WK-13", "WK-14", "Bridge"]):
        t.rows[0].cells[i].text = h
    for s in det:
        if s["delta"] == "N/A":
            continue
        row = t.add_row().cells
        row[0].text = s["site"]
        row[1].text = s["mp"]
        row[2].text = s["wk13"]
        row[3].text = f"{s['wk14']} ({s['delta']})"
        row[4].text = (s["bridge_wk14"] or "No bridge")[:120]

# Compliance by Site
doc.add_heading("WK-14 Compliance by Site", level=1)
add_img(doc, "chart_compliance_by_site.png")

# Root Cause Pareto
doc.add_heading("Root Cause Pareto", level=1)
add_img(doc, "chart_root_cause_pareto.png")

# Miss Type Breakdown
doc.add_heading("Miss Type Breakdown", level=1)
add_img(doc, "chart_miss_type_breakdown.png", 5)

# Top 5 Root Causes Detail
doc.add_heading("Top 5 Root Causes — Detail", level=1)
for rc in analysis.get("top_5_root_causes", []):
    doc.add_heading(f"#{rc['rank']}: {rc['root_cause']}", level=2)
    t = doc.add_table(rows=5, cols=2)
    t.style = "Light List Accent 1"
    for i, (label, val) in enumerate([
        ("Affected Sites", ", ".join(rc.get("affected_sites", []))),
        ("Impact", rc.get("impact", "")),
        ("Actions", "\n".join(f"• {a}" for a in rc.get("recommended_actions", []))),
        ("Owner", rc.get("owner", "TBD")),
        ("Timeline", rc.get("timeline", "TBD")),
    ]):
        t.rows[i].cells[0].text = label
        t.rows[i].cells[1].text = val

# Country Breakdown
doc.add_heading("Country Breakdown", level=1)
add_img(doc, "chart_by_country.png")

# New Flags
doc.add_heading("New Flags in WK-14", level=1)
nf = wow.get("new_flags", [])
if nf:
    t = doc.add_table(rows=1, cols=4)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Site", "MP", "WK-14", "Bridge"]):
        t.rows[0].cells[i].text = h
    for s in nf:
        row = t.add_row().cells
        row[0].text = s["site"]
        row[1].text = s["mp"]
        row[2].text = s["wk14"]
        row[3].text = (s["bridge_wk14"] or "No bridge")[:120]

# Resolved
doc.add_heading("Resolved from WK-13", level=1)
res = wow.get("resolved", [])
if res:
    t = doc.add_table(rows=1, cols=3)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Site", "MP", "Was (WK-13)"]):
        t.rows[0].cells[i].text = h
    for s in res:
        row = t.add_row().cells
        row[0].text = s["site"]
        row[1].text = s["mp"]
        row[2].text = s["wk13"]

# Site Detail Table
doc.add_heading("Full Site-Level Detail", level=1)
sites = analysis.get("site_summaries", [])
if sites:
    t = doc.add_table(rows=1, cols=5)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Site", "MP", "Compliance", "Bridge Summary", "Severity"]):
        t.rows[0].cells[i].text = h
    for s in sites:
        row = t.add_row().cells
        row[0].text = s.get("site", "")
        row[1].text = s.get("mp", "")
        row[2].text = str(s.get("compliance", ""))
        row[3].text = s.get("bridge_summary", "")
        row[4].text = s.get("severity", "")

# WBR Talking Points
doc.add_heading("WBR Talking Points", level=1)
for pt in analysis.get("recommended_wbr_talking_points", []):
    doc.add_paragraph(pt, style="List Bullet")

# Appendix
doc.add_heading("Appendix: Sites Without Bridges", level=1)
no_bridge = [s for s in sites if "No bridge" in s.get("bridge_summary", "")]
if no_bridge:
    t = doc.add_table(rows=1, cols=3)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Site", "MP", "Compliance"]):
        t.rows[0].cells[i].text = h
    for s in no_bridge:
        row = t.add_row().cells
        row[0].text = s["site"]
        row[1].text = s["mp"]
        row[2].text = s["compliance"]

word_path = f"{OUTPUT}/WBR_Bridge_Analysis_FINAL.docx"
doc.save(word_path)
print(f"Word: {word_path}")

# ============ POWERPOINT ============
prs = Presentation()
prs.slide_width = PptxInches(13.333)
prs.slide_height = PptxInches(7.5)


def add_slide_with_image(title, image_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.3), PptxInches(0.1), PptxInches(12), PptxInches(0.5))
    txBox.text_frame.paragraphs[0].text = title
    txBox.text_frame.paragraphs[0].font.size = PptxInches(0.25)
    txBox.text_frame.paragraphs[0].font.bold = True
    p = img(image_name)
    if os.path.exists(p):
        slide.shapes.add_picture(p, PptxInches(0.5), PptxInches(0.7), width=PptxInches(12), height=PptxInches(6.5))


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "EU LM WBR — CRET Scan Compliance Bridge"
slide.placeholders[1].text = (
    f"Week 14 | {analysis['total_flagged_sites']} Sites Flagged | "
    f"{len(wow['deteriorating'])} Deteriorating | {len(wow['resolved'])} Resolved"
)

# Slide 2: Executive Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Executive Summary"
slide.placeholders[1].text = (
    f"{analysis['executive_summary']}\n\n"
    f"WoW: {wow['summary']}"
)

# Slide 3-7: Charts
add_slide_with_image("Compliance Trend — WK-12 → WK-14", "chart_wow_trending.png")
add_slide_with_image("Deteriorating Sites: WK-13 → WK-14", "chart_deteriorating.png")
add_slide_with_image("WK-14 Compliance by Site", "chart_compliance_by_site.png")
add_slide_with_image("Root Cause Pareto", "chart_root_cause_pareto.png")
add_slide_with_image("Flagged Sites per Week by Country", "chart_wow_site_count.png")

# Slide 8: Top 5 Root Causes
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Top 5 Root Causes & Actions"
body = ""
for rc in analysis.get("top_5_root_causes", []):
    actions = rc.get("recommended_actions", [])
    body += (f"#{rc['rank']} {rc['root_cause']}\n"
             f"   Sites: {', '.join(rc.get('affected_sites', []))}\n"
             f"   Action: {actions[0] if actions else 'TBD'}\n"
             f"   Owner: {rc.get('owner', 'TBD')} | {rc.get('timeline', 'TBD')}\n\n")
slide.placeholders[1].text = body.strip()

# Slide 9: WBR Talking Points
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "WBR Talking Points"
slide.placeholders[1].text = "\n\n".join(f"• {pt}" for pt in analysis.get("recommended_wbr_talking_points", []))

pptx_path = f"{OUTPUT}/WBR_Bridge_Analysis_FINAL.pptx"
prs.save(pptx_path)
print(f"PPTX: {pptx_path}")
