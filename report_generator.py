"""Generate Word and PowerPoint outputs for WBR bridge analysis."""
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
import os


def _severity_color(severity: str) -> RGBColor:
    return {"HIGH": RGBColor(0xFF, 0x00, 0x00), "MEDIUM": RGBColor(0xFF, 0xA5, 0x00)}.get(
        severity.upper(), RGBColor(0x00, 0x80, 0x00)
    )


def generate_word(analysis: dict, output_path: str) -> str:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Title
    title = doc.add_heading("EU LM WBR — CRET Scan Compliance Bridge", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Executive Summary
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(analysis.get("executive_summary", "N/A"))
    doc.add_paragraph(f"Total Flagged Sites: {analysis.get('total_flagged_sites', 'N/A')}")

    # Top 5 Root Causes
    doc.add_heading("Top 5 Root Causes & Actions", level=1)
    for rc in analysis.get("top_5_root_causes", []):
        doc.add_heading(f"#{rc['rank']}: {rc['root_cause']}", level=2)
        doc.add_paragraph(f"Affected Sites: {', '.join(rc.get('affected_sites', []))}")
        doc.add_paragraph(f"Impact: {rc.get('impact', 'N/A')}")
        doc.add_paragraph("Recommended Actions:")
        for action in rc.get("recommended_actions", []):
            doc.add_paragraph(action, style="List Bullet")
        doc.add_paragraph(f"Owner: {rc.get('owner', 'TBD')} | Timeline: {rc.get('timeline', 'TBD')}")

    # Site-Level Summary Table
    doc.add_heading("Site-Level Summary", level=1)
    sites = analysis.get("site_summaries", [])
    if sites:
        table = doc.add_table(rows=1, cols=5)
        table.style = "Light Grid Accent 1"
        for i, header in enumerate(["Site", "MP", "Compliance", "Bridge Summary", "Severity"]):
            table.rows[0].cells[i].text = header

        for s in sites:
            row = table.add_row().cells
            row[0].text = s.get("site", "")
            row[1].text = s.get("mp", "")
            row[2].text = str(s.get("compliance", ""))
            row[3].text = s.get("bridge_summary", "")
            row[4].text = s.get("severity", "")

    # Patterns
    doc.add_heading("Patterns & Trends", level=1)
    doc.add_paragraph(analysis.get("patterns_and_trends", "N/A"))

    # WBR Talking Points
    doc.add_heading("Recommended WBR Talking Points", level=1)
    for point in analysis.get("recommended_wbr_talking_points", []):
        doc.add_paragraph(point, style="List Bullet")

    doc.save(output_path)
    return output_path


def generate_pptx(analysis: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EU LM WBR — CRET Scan Compliance Bridge"
    slide.placeholders[1].text = f"Flagged Sites: {analysis.get('total_flagged_sites', 'N/A')}"

    # Executive Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = analysis.get("executive_summary", "N/A")

    # Top 5 Root Causes slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top 5 Root Causes"
    body = ""
    for rc in analysis.get("top_5_root_causes", []):
        sites = ", ".join(rc.get("affected_sites", []))
        actions = "; ".join(rc.get("recommended_actions", []))
        body += f"#{rc['rank']} {rc['root_cause']}\n"
        body += f"   Sites: {sites} | Impact: {rc.get('impact', '')}\n"
        body += f"   Actions: {actions}\n\n"
    slide.placeholders[1].text = body.strip()

    # Site Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.title.text = "Site-Level Summary"
    sites = analysis.get("site_summaries", [])
    if sites:
        rows = len(sites) + 1
        cols = 5
        left = PptxInches(0.5)
        top = PptxInches(1.5)
        width = PptxInches(12)
        height = PptxInches(0.4 * rows)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        for i, header in enumerate(["Site", "MP", "Compliance", "Bridge", "Severity"]):
            table.cell(0, i).text = header

        for r, s in enumerate(sites, 1):
            table.cell(r, 0).text = s.get("site", "")
            table.cell(r, 1).text = s.get("mp", "")
            table.cell(r, 2).text = str(s.get("compliance", ""))
            table.cell(r, 3).text = s.get("bridge_summary", "")[:60]
            table.cell(r, 4).text = s.get("severity", "")

    prs.save(output_path)
    return output_path
