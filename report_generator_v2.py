"""Enhanced report generator with charts and WBR best practices."""
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
import os


def _add_chart_to_doc(doc, image_path, width_inches=6.5):
    if os.path.exists(image_path):
        doc.add_picture(image_path, width=Inches(width_inches))
        last_paragraph = doc.paragraphs[-1]
        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER


def generate_word_enhanced(analysis: dict, charts_dir: str, output_path: str) -> str:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Amazon Ember"
    style.font.size = Pt(10)
    style.paragraph_format.space_after = Pt(4)

    # === TITLE ===
    title = doc.add_heading("EU LM WBR — CRET Scan Compliance Bridge", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("Week 14 | Auto-generated Bridge Analysis", style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER

    # === WBR FORMAT: Start with data, then narrative ===

    # 1. Key Metrics Box
    doc.add_heading("Key Metrics", level=1)
    t = doc.add_table(rows=2, cols=4)
    t.style = "Light Grid Accent 1"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(["Total Flagged", "With Bridge", "Without Bridge", "Worst Site"]):
        t.rows[0].cells[i].text = h
    t.rows[1].cells[0].text = str(analysis.get("total_flagged_sites", ""))
    t.rows[1].cells[1].text = str(analysis.get("sites_with_bridges", ""))
    t.rows[1].cells[2].text = str(analysis.get("sites_without_bridges", ""))
    # Find worst site
    worst = min(analysis.get("site_summaries", []),
                key=lambda s: float(s["compliance"].strip("%")) if "%" in s["compliance"] else 100,
                default={"site": "N/A", "compliance": "N/A"})
    t.rows[1].cells[3].text = f"{worst['site']} ({worst['compliance']})"

    # 2. Compliance Chart
    doc.add_heading("Site Compliance Overview", level=1)
    _add_chart_to_doc(doc, os.path.join(charts_dir, "chart_compliance_by_site.png"))

    # 3. Executive Summary (WBR: concise, data-first)
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(analysis.get("executive_summary", "N/A"))

    # 4. Root Cause Pareto
    doc.add_heading("Root Cause Analysis", level=1)
    _add_chart_to_doc(doc, os.path.join(charts_dir, "chart_root_cause_pareto.png"))

    # 5. Miss Type Breakdown
    doc.add_heading("Miss Type Breakdown", level=1)
    _add_chart_to_doc(doc, os.path.join(charts_dir, "chart_miss_type_breakdown.png"))

    # 6. Top 5 Root Causes — Detail (WBR: RC → Impact → Action → Owner → Timeline)
    doc.add_heading("Top 5 Root Causes — Detail", level=1)
    for rc in analysis.get("top_5_root_causes", []):
        doc.add_heading(f"#{rc['rank']}: {rc['root_cause']}", level=2)

        # Use a mini table for structured RC detail
        t = doc.add_table(rows=5, cols=2)
        t.style = "Light List Accent 1"
        fields = [
            ("Affected Sites", ", ".join(rc.get("affected_sites", []))),
            ("Impact", rc.get("impact", "N/A")),
            ("Actions", "\n".join(f"• {a}" for a in rc.get("recommended_actions", []))),
            ("Owner", rc.get("owner", "TBD")),
            ("Timeline", rc.get("timeline", "TBD")),
        ]
        for i, (label, val) in enumerate(fields):
            t.rows[i].cells[0].text = label
            t.rows[i].cells[0].paragraphs[0].runs[0].bold = True if t.rows[i].cells[0].paragraphs[0].runs else None
            t.rows[i].cells[1].text = val

    # 7. Country Breakdown
    doc.add_heading("Country Breakdown", level=1)
    _add_chart_to_doc(doc, os.path.join(charts_dir, "chart_by_country.png"))

    # 8. Site-Level Detail Table
    doc.add_heading("Site-Level Detail", level=1)
    sites = analysis.get("site_summaries", [])
    if sites:
        table = doc.add_table(rows=1, cols=5)
        table.style = "Light Grid Accent 1"
        for i, header in enumerate(["Site", "MP", "Compliance", "Bridge Summary", "Severity"]):
            cell = table.rows[0].cells[i]
            cell.text = header
            for p in cell.paragraphs:
                for r in p.runs:
                    r.bold = True

        for s in sites:
            row = table.add_row().cells
            row[0].text = s.get("site", "")
            row[1].text = s.get("mp", "")
            row[2].text = str(s.get("compliance", ""))
            row[3].text = s.get("bridge_summary", "")
            row[4].text = s.get("severity", "")

    # 9. Patterns & Trends
    doc.add_heading("Patterns & Trends", level=1)
    for line in analysis.get("patterns_and_trends", "").split("\n"):
        if line.strip():
            doc.add_paragraph(line.strip())

    # 10. WBR Talking Points (WBR: what to say in the meeting)
    doc.add_heading("Recommended WBR Talking Points", level=1)
    for point in analysis.get("recommended_wbr_talking_points", []):
        doc.add_paragraph(point, style="List Bullet")

    # 11. Appendix: Missing Bridges
    doc.add_heading("Appendix: Sites Without Bridges — Action Required", level=1)
    no_bridge = [s for s in sites if "No bridge" in s.get("bridge_summary", "")]
    if no_bridge:
        t = doc.add_table(rows=1, cols=3)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "Compliance"]):
            t.rows[0].cells[i].text = h
        for s in no_bridge:
            row = t.add_row().cells
            row[0].text = s["site"]
            row[1].text = s["mp"]
            row[2].text = s["compliance"]

    doc.save(output_path)
    return output_path


def generate_pptx_enhanced(analysis: dict, charts_dir: str, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EU LM WBR — CRET Scan Compliance Bridge"
    slide.placeholders[1].text = (
        f"Week 14 | {analysis.get('total_flagged_sites', '?')} Sites Flagged | "
        f"{analysis.get('sites_with_bridges', '?')} Bridges Provided"
    )

    # Slide 2: Executive Summary + Key Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = analysis.get("executive_summary", "")

    # Slide 3: Compliance Chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    chart_path = os.path.join(charts_dir, "chart_compliance_by_site.png")
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, PptxInches(0.5), PptxInches(0.3),
                                  width=PptxInches(12), height=PptxInches(6.8))

    # Slide 4: Root Cause Pareto
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_path = os.path.join(charts_dir, "chart_root_cause_pareto.png")
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, PptxInches(1), PptxInches(0.5),
                                  width=PptxInches(11), height=PptxInches(6))

    # Slide 5: Miss Type Breakdown
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_path = os.path.join(charts_dir, "chart_miss_type_breakdown.png")
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, PptxInches(2.5), PptxInches(0.5),
                                  width=PptxInches(8), height=PptxInches(6))

    # Slide 6: Top 5 Root Causes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top 5 Root Causes & Actions"
    body = ""
    for rc in analysis.get("top_5_root_causes", []):
        sites = ", ".join(rc.get("affected_sites", []))
        actions = rc.get("recommended_actions", [])
        body += f"#{rc['rank']} {rc['root_cause']}\n"
        body += f"   Sites: {sites}\n"
        body += f"   Impact: {rc.get('impact', '')[:120]}\n"
        body += f"   Action: {actions[0] if actions else 'TBD'}\n"
        body += f"   Owner: {rc.get('owner', 'TBD')} | Timeline: {rc.get('timeline', 'TBD')}\n\n"
    slide.placeholders[1].text = body.strip()

    # Slide 7: Country Breakdown
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_path = os.path.join(charts_dir, "chart_by_country.png")
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, PptxInches(2), PptxInches(0.5),
                                  width=PptxInches(9), height=PptxInches(6))

    # Slide 8: WBR Talking Points
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "WBR Talking Points"
    slide.placeholders[1].text = "\n\n".join(
        f"• {pt}" for pt in analysis.get("recommended_wbr_talking_points", [])
    )

    prs.save(output_path)
    return output_path
