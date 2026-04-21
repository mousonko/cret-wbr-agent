"""Generate WBR report with narrative-style root cause bridges."""
import json
import os
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches


def img(output_dir, name):
    return os.path.join(output_dir, name)


def add_img(doc, output_dir, name, width=6.2):
    p = img(output_dir, name)
    if os.path.exists(p):
        doc.add_picture(p, width=Inches(width))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER


def _add_bold_run(paragraph, text):
    run = paragraph.add_run(text)
    run.bold = True
    return run


def _add_run(paragraph, text):
    return paragraph.add_run(text)


def generate_narrative_word(analysis, wow, output_dir):
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10)

    doc.add_heading("EU LM WBR — CRET Scan Compliance Bridge", level=0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    week = analysis.get("week", "")
    doc.add_paragraph(f"{week} | Comprehensive Analysis with Week-over-Week Trends",
                      style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Key Metrics
    doc.add_heading("Key Metrics", level=1)
    t = doc.add_table(rows=2, cols=6)
    t.style = "Light Grid Accent 1"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(["Flagged Sites", "With Bridge", "No Bridge", "Deteriorating", "New Flags", "Resolved"]):
        t.rows[0].cells[i].text = h
    t.rows[1].cells[0].text = str(analysis.get("total_flagged_sites", ""))
    t.rows[1].cells[1].text = str(analysis.get("sites_with_bridges", ""))
    t.rows[1].cells[2].text = str(analysis.get("sites_without_bridges", ""))
    t.rows[1].cells[3].text = str(len(wow.get("deteriorating", [])))
    t.rows[1].cells[4].text = str(len(wow.get("new_flags", [])))
    t.rows[1].cells[5].text = str(len(wow.get("resolved", [])))

    # Charts
    doc.add_heading("Compliance Trend", level=1)
    add_img(doc, output_dir, "chart_wow_trending.png")
    doc.add_heading("Flagged Sites per Week by Country", level=1)
    add_img(doc, output_dir, "chart_wow_site_count.png")
    doc.add_heading("Deteriorating Sites", level=1)
    add_img(doc, output_dir, "chart_deteriorating.png", 6.5)
    doc.add_heading("Compliance by Site", level=1)
    add_img(doc, output_dir, "chart_compliance_by_site.png")
    doc.add_heading("Root Cause Pareto", level=1)
    add_img(doc, output_dir, "chart_root_cause_pareto.png")
    doc.add_heading("Miss Type Breakdown", level=1)
    add_img(doc, output_dir, "chart_miss_type_breakdown.png", 5)

    # === NARRATIVE BRIDGE SECTION ===
    doc.add_heading("Insights Bridge Compliance", level=1)

    # Bridge summary paragraph
    bridge_summary = analysis.get("bridge_summary", "")
    if bridge_summary:
        p = doc.add_paragraph()
        _add_run(p, bridge_summary)

    # Top root causes as narrative bullets
    doc.add_paragraph()  # spacing
    num_causes = len(analysis.get("top_root_causes_narrative", []))
    p = doc.add_paragraph()
    _add_bold_run(p, f"Top {num_causes} root causes:")

    for rc in analysis.get("top_root_causes_narrative", []):
        p = doc.add_paragraph(style="List Bullet")
        _add_bold_run(p, f"RC{rc['rank']} {rc['title']}")
        _add_run(p, f" {rc['narrative']}")

    # Country Breakdown
    doc.add_heading("Country Breakdown", level=1)
    add_img(doc, output_dir, "chart_by_country.png")

    # Deteriorating table
    det = wow.get("deteriorating", [])
    if det:
        doc.add_heading("Deteriorating Sites Detail", level=1)
        t = doc.add_table(rows=1, cols=5)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "WK-13", "WK-14", "Bridge"]):
            t.rows[0].cells[i].text = h
        for s in det:
            if s["delta"] == "N/A":
                continue
            row = t.add_row().cells
            row[0].text = s["site"]
            row[1].text = s["mp"]
            row[2].text = s["wk13"]
            row[3].text = f"{s['wk14']} ({s['delta']})"
            row[4].text = (s["bridge_wk14"] or "No bridge")[:120]

    # New Flags
    nf = wow.get("new_flags", [])
    if nf:
        doc.add_heading("New Flags", level=1)
        t = doc.add_table(rows=1, cols=4)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "Compliance", "Bridge"]):
            t.rows[0].cells[i].text = h
        for s in nf:
            row = t.add_row().cells
            row[0].text = s["site"]
            row[1].text = s["mp"]
            row[2].text = s["wk14"]
            row[3].text = (s["bridge_wk14"] or "No bridge")[:120]

    # Resolved
    res = wow.get("resolved", [])
    if res:
        doc.add_heading("Resolved from Previous Week", level=1)
        t = doc.add_table(rows=1, cols=3)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "Was"]):
            t.rows[0].cells[i].text = h
        for s in res:
            row = t.add_row().cells
            row[0].text = s["site"]
            row[1].text = s["mp"]
            row[2].text = s["wk13"]

    # Full site table
    sites = analysis.get("site_summaries", [])
    if sites:
        doc.add_heading("Full Site-Level Detail", level=1)
        t = doc.add_table(rows=1, cols=5)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "Compliance", "Bridge Summary", "Severity"]):
            t.rows[0].cells[i].text = h
        for s in sites:
            row = t.add_row().cells
            row[0].text = s.get("site", "")
            row[1].text = s.get("mp", "")
            row[2].text = str(s.get("compliance", ""))
            row[3].text = s.get("bridge_summary", "")
            row[4].text = s.get("severity", "")

    # Appendix
    no_bridge = [s for s in sites if "No bridge" in s.get("bridge_summary", "")]
    if no_bridge:
        doc.add_heading("Appendix: Sites Without Bridges", level=1)
        t = doc.add_table(rows=1, cols=3)
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(["Site", "MP", "Compliance"]):
            t.rows[0].cells[i].text = h
        for s in no_bridge:
            row = t.add_row().cells
            row[0].text = s["site"]
            row[1].text = s["mp"]
            row[2].text = s["compliance"]

    word_path = os.path.join(output_dir, "WBR_Bridge_Analysis_FINAL.docx")
    doc.save(word_path)
    print(f"Word: {word_path}")
    return word_path


def generate_narrative_pptx(analysis, wow, output_dir):
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_slide_img(title, image_name):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.3), PptxInches(0.1), PptxInches(12), PptxInches(0.5))
        txBox.text_frame.paragraphs[0].text = title
        txBox.text_frame.paragraphs[0].font.size = PptxInches(0.25)
        txBox.text_frame.paragraphs[0].font.bold = True
        p = img(output_dir, image_name)
        if os.path.exists(p):
            slide.shapes.add_picture(p, PptxInches(0.5), PptxInches(0.7), width=PptxInches(12), height=PptxInches(6.5))

    week = analysis.get("week", "")

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EU LM WBR — CRET Scan Compliance Bridge"
    slide.placeholders[1].text = (
        f"{week} | {analysis['total_flagged_sites']} Sites Flagged | "
        f"{len(wow.get('deteriorating', []))} Deteriorating | {len(wow.get('resolved', []))} Resolved"
    )

    # Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = analysis.get("bridge_summary", analysis.get("executive_summary", ""))

    # Charts
    add_slide_img("Compliance Trend", "chart_wow_trending.png")
    add_slide_img("Deteriorating Sites", "chart_deteriorating.png")
    add_slide_img("Compliance by Site", "chart_compliance_by_site.png")
    add_slide_img("Root Cause Pareto", "chart_root_cause_pareto.png")
    add_slide_img("Flagged Sites per Week by Country", "chart_wow_site_count.png")

    # Narrative Root Causes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Insights Bridge Compliance"
    body = ""
    for rc in analysis.get("top_root_causes_narrative", []):
        body += f"RC{rc['rank']} {rc['title']}: {rc['narrative']}\n\n"
    slide.placeholders[1].text = body.strip()

    # WBR Talking Points
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "WBR Talking Points"
    slide.placeholders[1].text = "\n\n".join(
        f"• {pt}" for pt in analysis.get("recommended_wbr_talking_points", [])
    )

    pptx_path = os.path.join(output_dir, "WBR_Bridge_Analysis_FINAL.pptx")
    prs.save(pptx_path)
    print(f"PPTX: {pptx_path}")
    return pptx_path
